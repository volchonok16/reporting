from __future__ import annotations

import io
import logging
import re
from copy import deepcopy
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from zoneinfo import ZoneInfo

import httpx
from fastapi import HTTPException
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

from app.b2b_news_service import load_b2b_news
from app.config import settings
from app.product_status_slides_template import fetch_google_slides_pptx
from app.product_status_rich_text import (
    CellStyle,
    display_cell_text,
    split_cell_wrapper,
    split_style_segments,
)
from app.product_status_service import load_b2b_product_status
from app.schemas import ProductStatusB2BOut, ProductStatusSheetOut

logger = logging.getLogger(__name__)

MOSCOW_TZ = ZoneInfo("Europe/Moscow")
DATE_PLACEHOLDER = "<%date>"
MARKET_NEWS_SLIDE_INDEX = 1
MARKET_NEWS_MARKERS = ("<$date>", "<$news>", "<$description>")
FIXED_SLIDE_COUNT = 3
COLUMN_COUNT = 4
WHY_COLUMN_INDEX = 3
DESCRIPTION_SLOT_INDEX = 2
DESCRIPTION_HEADER = "Описание проекта и статус"
WHY_HEADER = "Зачем и для чего делаем"
PRESENTATION_FLAG_PATTERN = r"идет в презентацию"
ATTENTION_FLAG_PATTERN = r"обратить.*вним"
PRESENTATION_DESCRIPTION_PATTERN = r"для презентации"
FULL_DESCRIPTION_PATTERN = r"полное описание"
WHY_PATTERN = r"зачем"
DESCRIPTION_PATTERN = r"описан"

TABLE_FONT_NAME = "T2 Rooftop"
TITLE_FONT_NAME = "T2 Halvar Breit ExtraBold"
TABLE_FONT_SIZE = Pt(10)
TABLE_FONT_SIZE_DENSE = Pt(8)
TITLE_FONT_SIZE = Pt(22)
TITLE_COLOR = RGBColor(0x00, 0x00, 0x00)
TEXT_COLOR = RGBColor(0x00, 0x00, 0x00)
ATTENTION_TEXT_COLOR = RGBColor(0xC0, 0x00, 0x00)
WHITE_FILL = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_BORDER_COLOR = "7F7F7F"
TABLE_BORDER_WIDTH = "12700"
DEFAULT_HIGHLIGHT_COLOR = "FFFF00"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

CHAR_WIDTH_EMU = 76000
LINE_HEIGHT_EMU = 145000
LINE_HEIGHT_DENSE_EMU = int(LINE_HEIGHT_EMU * TABLE_FONT_SIZE_DENSE.pt / TABLE_FONT_SIZE.pt)
CELL_TEXT_MARGIN_EMU = int(Pt(6).emu)
HEADER_ROW_HEIGHT_EMU = 320000
MIN_ROW_HEIGHT_EMU = LINE_HEIGHT_EMU + CELL_TEXT_MARGIN_EMU
BODY_HEIGHT_EMU = 3_408_600
BODY_WIDTH_EMU = 8_900_700
COLUMN_WIDTH_RATIOS = (0.10, 0.14, 0.46, 0.30)
NEWS_BODY_MARGIN_LEFT_EMU = 457200
NEWS_CHAR_WIDTH_EMU = 55000
NEWS_LINE_HEIGHT_EMU = 152400
NEWS_PARAGRAPH_GAP_EMU = 40000
NOTES_BLOCK_SEPARATOR = "—"
NOTES_MUTED_COLOR = RGBColor(0xA0, 0xA0, 0xA0)
NOTES_BLOCK_SPACE_BEFORE = Pt(12)
NOTES_PARAGRAPH_SPACE_AFTER = Pt(6)
CELL_PARAGRAPH_SPACE_BEFORE = Pt(4)
CELL_PARAGRAPH_SPACE_AFTER = Pt(3)

_CNVPR_TAGS = (
    "{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr",
    "{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr",
)

# 0-based индексы слайдов-образцов в Status.pptx (слайд 4 = индекс 3).
# Ключ — часть названия листа после «:» (префикс «Офис» / «Продуктовый офис» не важен).
_SHEET_TEMPLATE_INDEX: dict[str, int] = {
    "core": 3,
    "m2m / iot": 4,
    "sms": 5,
    "voice": 6,
    "перспективные продукты": 7,
    "продуктовый маркетинг": 3,
}

# Совместимость со старыми тестами/импортами.
COVER_SECTION_NAME = "Раздел по умолчанию"
TABLE_FONT_SZ_XML = "1000"
TABLE_FONT_SZ_DENSE_XML = "800"
_CELL_FONT_BY_INDEX: tuple[tuple[Pt, str], ...] = (
    (TABLE_FONT_SIZE, TABLE_FONT_SZ_XML),
    (TABLE_FONT_SIZE, TABLE_FONT_SZ_XML),
    (TABLE_FONT_SIZE, TABLE_FONT_SZ_XML),
    (TABLE_FONT_SIZE_DENSE, TABLE_FONT_SZ_DENSE_XML),
)


@dataclass(frozen=True)
class ContentSlideTemplate:
    index: int
    title: str = ""
    row_count: int = 0
    column_count: int = COLUMN_COUNT
    table_height: int = BODY_HEIGHT_EMU
    col_chars_per_line: tuple[int, ...] = (12, 16, 54, 34)


class TemplateCatalog:
    """Минимальная обёртка для совместимости тестов — шаблон Status.pptx."""

    def __init__(self, *, title_slide_index: int, templates: list[ContentSlideTemplate]) -> None:
        self.title_slide_index = title_slide_index
        self._templates = templates
        self._by_title: dict[str, list[ContentSlideTemplate]] = {}
        for template in templates:
            key = _normalize_title(template.title)
            self._by_title.setdefault(key, []).append(template)
        self._default = templates[0] if templates else ContentSlideTemplate(index=3)

    @classmethod
    def from_presentation(cls, prs: Presentation) -> TemplateCatalog:
        templates = [
            ContentSlideTemplate(index=index, title=_slide_title(prs.slides[index]) or "")
            for index in sorted(set(_SHEET_TEMPLATE_INDEX.values()))
            if index < len(prs.slides)
        ]
        return cls(title_slide_index=0, templates=templates)

    def match(self, sheet_name: str) -> list[ContentSlideTemplate]:
        index = _template_index_for_sheet(sheet_name)
        return [ContentSlideTemplate(index=index, title=sheet_name)]

    def template_pool(self, sheet_name: str) -> list[ContentSlideTemplate]:
        return self.match(sheet_name)

    def chunk_plan(
        self,
        sheet_name: str,
        rows: list[dict[str, str]],
        columns: list[str],
    ) -> list[tuple[ContentSlideTemplate, list[dict[str, str]]]]:
        template = self.match(sheet_name)[0]
        chunks = _chunk_rows(rows, columns)
        if not chunks:
            return [(template, [])]
        return [(template, chunk) for chunk in chunks]


def _normalize_title(value: str) -> str:
    return re.sub(r"\s+", " ", value.strip().casefold())


def _sheet_section_key(sheet_name: str) -> str:
    """Нормализованная часть названия листа после двоеточия."""
    short = sheet_name.split(":", 1)[1].strip() if ":" in sheet_name else sheet_name.strip()
    return _normalize_title(short)


def _looks_like_url(value: str) -> bool:
    lowered = value.strip().casefold()
    return lowered.startswith("http://") or lowered.startswith("https://")


def _bundled_template_path() -> Path:
    return Path(__file__).resolve().parent.parent / "assets" / "b2b_product_status_template.pptx"


def _google_http_client() -> httpx.Client:
    proxy = settings.outbound_http_proxy.strip() or None
    return httpx.Client(timeout=60.0, follow_redirects=True, proxy=proxy)


def _open_template_presentation() -> Presentation:
    configured = settings.b2b_product_status_presentation_template.strip()
    if configured:
        if _looks_like_url(configured):
            logger.warning(
                "product_status_template_url_in_template_setting — "
                "используйте B2B_PRODUCT_STATUS_PRESENTATION_REFERENCE_URL"
            )
        else:
            path = Path(configured)
            if path.is_file():
                return Presentation(str(path))
            logger.warning("product_status_template_missing path=%s", configured)

    reference_url = settings.b2b_product_status_presentation_reference_url.strip()
    if reference_url and _looks_like_url(reference_url):
        try:
            with _google_http_client() as client:
                pptx_bytes = fetch_google_slides_pptx(reference_url=reference_url, client=client)
            prs = Presentation(io.BytesIO(pptx_bytes))
            if len(prs.slides) >= FIXED_SLIDE_COUNT + 1:
                return prs
            logger.warning(
                "google_slides_template_too_short slides=%s url=%s",
                len(prs.slides),
                reference_url,
            )
        except (httpx.HTTPError, ValueError, OSError):
            logger.warning(
                "google_slides_template_fetch_failed url=%s",
                reference_url,
                exc_info=True,
            )

    default = _bundled_template_path()
    if default.is_file():
        return Presentation(str(default))

    raise HTTPException(
        status_code=503,
        detail=(
            "Шаблон презентации не найден. Проверьте доступ к эталону Google Slides "
            "или положите PPTX в backend/assets/b2b_product_status_template.pptx."
        ),
    )


def _template_path() -> Path:
    """Совместимость для тестов и скриптов — предпочтительно _open_template_presentation()."""
    configured = settings.b2b_product_status_presentation_template.strip()
    if configured and not _looks_like_url(configured):
        path = Path(configured)
        if path.is_file():
            return path
    default = _bundled_template_path()
    if default.is_file():
        return default
    raise HTTPException(status_code=503, detail="Шаблон презентации статуса продукта B2B не найден.")


def _template_index_for_sheet(sheet_name: str) -> int:
    return _SHEET_TEMPLATE_INDEX.get(_sheet_section_key(sheet_name), 3)


def _section_name_for_sheet(sheet_name: str) -> str:
    short = sheet_name.split(":", 1)[1].strip() if ":" in sheet_name else sheet_name.strip()
    aliases = {
        "voice": "Voice",
        "продуктовый маркетинг": "Продуктовый Маркетинг",
    }
    return aliases.get(short.casefold(), short)


def _find_column(
    columns: list[str],
    pattern: str,
    *,
    exclude: set[str] | None = None,
) -> str:
    regex = re.compile(pattern, re.IGNORECASE)
    excluded = exclude or set()
    for column in columns:
        if column in excluded:
            continue
        if regex.search(column.strip()):
            return column
    return ""


def _is_description_presentation_column(column: str) -> bool:
    key = column.strip().casefold()
    return (
        re.search(PRESENTATION_DESCRIPTION_PATTERN, key) is not None
        and re.search(DESCRIPTION_PATTERN, key) is not None
        and re.search(WHY_PATTERN, key) is None
    )


def _is_why_presentation_column(column: str) -> bool:
    key = column.strip().casefold()
    return (
        re.search(PRESENTATION_DESCRIPTION_PATTERN, key) is not None
        and re.search(WHY_PATTERN, key) is not None
    )


def _is_full_description_notes_column(column: str) -> bool:
    key = column.strip().casefold()
    return (
        re.search(FULL_DESCRIPTION_PATTERN, key) is not None
        and re.search(DESCRIPTION_PATTERN, key) is not None
        and re.search(WHY_PATTERN, key) is None
    )


def _is_full_why_notes_column(column: str) -> bool:
    key = column.strip().casefold()
    return (
        re.search(FULL_DESCRIPTION_PATTERN, key) is not None
        and re.search(WHY_PATTERN, key) is not None
    )


def _is_presentation_internal_column(column: str) -> bool:
    key = column.strip().casefold()
    if key == "зни":
        return True
    if re.search(PRESENTATION_FLAG_PATTERN, key):
        return True
    if re.search(r"обратить.*вним", key):
        return True
    if re.search(r"комментар", key):
        return True
    if _is_full_description_notes_column(column):
        return True
    if _is_full_why_notes_column(column):
        return True
    return False


def _is_legacy_description_column(column: str) -> bool:
    key = column.strip().casefold()
    return (
        re.search(DESCRIPTION_PATTERN, key) is not None
        and re.search(WHY_PATTERN, key) is None
        and not re.search(PRESENTATION_DESCRIPTION_PATTERN, key)
        and not re.search(FULL_DESCRIPTION_PATTERN, key)
        and not key.startswith("полное")
    )


def _is_legacy_why_column(column: str) -> bool:
    key = column.strip().casefold()
    return (
        re.search(WHY_PATTERN, key) is not None
        and not re.search(PRESENTATION_DESCRIPTION_PATTERN, key)
        and not re.search(FULL_DESCRIPTION_PATTERN, key)
    )


def _full_description_notes_column(columns: list[str]) -> str:
    for column in columns:
        if _is_full_description_notes_column(column):
            return column
    for column in columns:
        key = column.strip().casefold()
        if (
            key.startswith("полное")
            and re.search(DESCRIPTION_PATTERN, key) is not None
            and re.search(WHY_PATTERN, key) is None
            and not re.search(PRESENTATION_DESCRIPTION_PATTERN, key)
        ):
            return column
    return ""


def _legacy_description_notes_column(columns: list[str]) -> str:
    for column in columns:
        if _is_legacy_description_column(column):
            return column
    return ""


def _legacy_why_notes_column(columns: list[str]) -> str:
    for column in columns:
        if _is_legacy_why_column(column):
            return column
    return ""


def _description_notes_source_columns(columns: list[str]) -> list[str]:
    """Столбцы с полным описанием для заметок (от приоритетного к запасному)."""
    sources: list[str] = []
    full_column = _full_description_notes_column(columns)
    if full_column:
        sources.append(full_column)
    if any(_is_description_presentation_column(column) for column in columns):
        legacy = _legacy_description_notes_column(columns)
        if legacy and legacy not in sources:
            sources.append(legacy)
    elif not sources:
        table_column = _description_value_column(columns)
        if table_column and not _is_presentation_internal_column(table_column):
            sources.append(table_column)
    return sources


def _why_notes_source_columns(columns: list[str]) -> list[str]:
    """Столбцы с полным «зачем» для заметок (от приоритетного к запасному)."""
    sources: list[str] = []
    full_column = _full_why_notes_column(columns)
    if full_column:
        sources.append(full_column)
    if any(_is_why_presentation_column(column) for column in columns):
        legacy = _legacy_why_notes_column(columns)
        if legacy and legacy not in sources:
            sources.append(legacy)
    elif not sources:
        table_column = _why_value_column(columns)
        if table_column and not _is_presentation_internal_column(table_column):
            sources.append(table_column)
    return sources


def _notes_text_from_sources(
    row: dict[str, str],
    source_columns: list[str],
    *,
    presentation_column: str = "",
) -> str:
    presentation_text = (
        display_cell_text(row.get(presentation_column, "")).strip() if presentation_column else ""
    )
    for column in source_columns:
        text = display_cell_text(row.get(column, "")).strip()
        if not text:
            continue
        if presentation_column and text == presentation_text:
            continue
        return text
    return ""


def _full_why_notes_column(columns: list[str]) -> str:
    for column in columns:
        if _is_full_why_notes_column(column):
            return column
    for column in columns:
        key = column.strip().casefold()
        if (
            key.startswith("полное")
            and re.search(WHY_PATTERN, key) is not None
            and not re.search(PRESENTATION_DESCRIPTION_PATTERN, key)
        ):
            return column
    return ""


def _description_value_column(columns: list[str]) -> str:
    for column in columns:
        if _is_description_presentation_column(column):
            return column

    has_presentation_desc = any(_is_description_presentation_column(column) for column in columns)
    excluded = {column for column in columns if _is_presentation_internal_column(column)}
    excluded.update(column for column in columns if _is_description_presentation_column(column))
    excluded.update(column for column in columns if _is_why_presentation_column(column))
    for column in columns:
        if column in excluded:
            continue
        if has_presentation_desc and _is_legacy_description_column(column):
            continue
        key = column.strip().casefold()
        if re.search(DESCRIPTION_PATTERN, key) and not re.search(WHY_PATTERN, key):
            return column
    return ""


def _why_value_column(columns: list[str]) -> str:
    for column in columns:
        if _is_why_presentation_column(column):
            return column

    has_presentation_why = any(_is_why_presentation_column(column) for column in columns)
    excluded = {column for column in columns if _is_presentation_internal_column(column)}
    excluded.update(column for column in columns if _is_why_presentation_column(column))
    excluded.update(column for column in columns if _is_description_presentation_column(column))
    for column in columns:
        if column in excluded:
            continue
        if has_presentation_why and _is_legacy_why_column(column):
            continue
        if re.search(WHY_PATTERN, column.strip(), re.IGNORECASE):
            return column
    return ""


def _mapped_columns(columns: list[str]) -> list[str]:
    usable = [column.strip() for column in columns if column.strip()]
    mapped = [
        _find_column(columns, r"^Дата"),
        _find_column(columns, r"^Проект"),
        _description_value_column(columns),
        _why_value_column(columns),
    ]

    used = {column for column in mapped if column}
    remaining = [column for column in usable if column not in used]
    for index, column_name in enumerate(mapped):
        if column_name or not remaining:
            continue
        mapped[index] = remaining.pop(0)
        used.add(mapped[index])

    return mapped


def _presentation_headers(columns: list[str]) -> list[str]:
    mapped = _mapped_columns(columns)
    defaults = ("Дата запуска", "Проект", DESCRIPTION_HEADER, WHY_HEADER)
    headers: list[str] = []
    for index, (column_name, default_name) in enumerate(zip(mapped, defaults, strict=True)):
        if index == DESCRIPTION_SLOT_INDEX:
            headers.append(DESCRIPTION_HEADER)
        elif index == WHY_COLUMN_INDEX:
            headers.append(WHY_HEADER)
        else:
            headers.append(column_name or default_name)
    return headers


def _is_yes_cell_value(value: str) -> bool:
    normalized = display_cell_text(value or "").strip().casefold()
    if normalized in ("нет", "no", "0", "false"):
        return False
    return normalized in ("да", "yes", "1", "true")


def _attention_column(columns: list[str]) -> str:
    return _find_column(columns, ATTENTION_FLAG_PATTERN)


def _is_attention_row(row: dict[str, str], columns: list[str]) -> bool:
    column = _attention_column(columns)
    if not column:
        return False
    return _is_yes_cell_value(row.get(column, ""))


def _filter_presentation_rows(
    rows: list[dict[str, str]],
    columns: list[str],
) -> list[dict[str, str]]:
    flag_column = _find_column(columns, PRESENTATION_FLAG_PATTERN)
    if not flag_column:
        return list(rows)
    return [
        row
        for row in rows
        if _is_yes_cell_value(row.get(flag_column, ""))
    ]


def _description_paragraphs(text: str) -> list[str]:
    paragraphs: list[str] = []
    for chunk in re.split(r"\n\s*\n", text.strip()):
        stripped = chunk.strip()
        if stripped:
            paragraphs.append(stripped)
    if paragraphs:
        return paragraphs
    stripped = text.strip()
    return [stripped] if stripped else []


def _slide_notes_blocks(rows: list[dict[str, str]], columns: list[str]) -> list[list[str]]:
    description_sources = _description_notes_source_columns(columns)
    why_sources = _why_notes_source_columns(columns)
    project_column = _find_column(columns, r"^Проект")
    if not description_sources and not why_sources:
        return []

    presentation_description_column = (
        _description_value_column(columns)
        if any(_is_description_presentation_column(column) for column in columns)
        else ""
    )
    presentation_why_column = (
        _why_value_column(columns)
        if any(_is_why_presentation_column(column) for column in columns)
        else ""
    )

    blocks: list[list[str]] = []
    for row in rows:
        description_text = _notes_text_from_sources(
            row,
            description_sources,
            presentation_column=presentation_description_column,
        )
        why_text = _notes_text_from_sources(
            row,
            why_sources,
            presentation_column=presentation_why_column,
        )
        if not description_text and not why_text:
            continue
        project = (
            display_cell_text(row.get(project_column, "")).strip()
            if project_column
            else ""
        )
        block: list[str] = []
        if project:
            block.append(project)
        if description_text:
            block.append(description_text)
        if why_text:
            block.append(why_text)
        blocks.append(block)
    return blocks


def _slide_notes_text(rows: list[dict[str, str]], columns: list[str]) -> str:
    """Совместимость для тестов — плоское представление блоков заметок."""
    rendered: list[str] = []
    for block_index, block in enumerate(_slide_notes_blocks(rows, columns)):
        if block_index > 0:
            rendered.append(NOTES_BLOCK_SEPARATOR)
        rendered.append("\n\n".join(block))
    return "\n\n".join(rendered)


def _row_values(row: dict[str, str], columns: list[str], col_count: int) -> list[str]:
    mapped = _mapped_columns(columns)
    values: list[str] = []
    for index in range(col_count):
        column_name = mapped[index] if index < len(mapped) else ""
        values.append((row.get(column_name, "") if column_name else "").strip())
    return values


def _cell_font_size(col_index: int, text: str = "") -> tuple[Pt, str]:
    del text
    if 0 <= col_index < len(_CELL_FONT_BY_INDEX):
        return _CELL_FONT_BY_INDEX[col_index]
    return TABLE_FONT_SIZE, TABLE_FONT_SZ_XML


def _column_chars_per_line(widths: tuple[int, ...]) -> tuple[int, ...]:
    dense_scale = TABLE_FONT_SIZE.pt / TABLE_FONT_SIZE_DENSE.pt
    chars: list[int] = []
    for index, width in enumerate(widths):
        per_line = max(8, int(width / CHAR_WIDTH_EMU))
        if index == WHY_COLUMN_INDEX:
            per_line = max(8, int(per_line * dense_scale))
        chars.append(per_line)
    return tuple(chars)


def _column_widths(total_width: int) -> tuple[int, ...]:
    widths = [int(total_width * ratio) for ratio in COLUMN_WIDTH_RATIOS]
    widths[-1] = total_width - sum(widths[:-1])
    return tuple(widths)


def _sanitize_cell_text(text: str) -> str:
    cleaned = (text or "").replace("\x0b", "\n").replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in cleaned.split("\n")]
    return "\n".join(lines).strip()


def _estimate_text_lines(text: str, chars_per_line: int) -> int:
    cleaned = display_cell_text(text)
    if not cleaned:
        return 0
    total = 0
    for line in cleaned.split("\n"):
        stripped = line.strip()
        if not stripped:
            total += 1
            continue
        total += max(1, (len(stripped) + chars_per_line - 1) // chars_per_line)
    return total


def _estimate_cell_height(text: str, col_index: int, col_chars_per_line: tuple[int, ...]) -> int:
    chars_per_line = col_chars_per_line[col_index] if col_index < len(col_chars_per_line) else 40
    lines = _estimate_text_lines(text, chars_per_line)
    if not lines:
        return 0
    line_height = LINE_HEIGHT_DENSE_EMU if col_index == WHY_COLUMN_INDEX else LINE_HEIGHT_EMU
    return lines * line_height + CELL_TEXT_MARGIN_EMU


def _estimate_row_height(values: list[str], col_chars_per_line: tuple[int, ...]) -> int:
    heights = [
        _estimate_cell_height(text, index, col_chars_per_line) for index, text in enumerate(values)
    ]
    positive = [height for height in heights if height > 0]
    if not positive:
        return MIN_ROW_HEIGHT_EMU
    return max(MIN_ROW_HEIGHT_EMU, max(positive))


def _chunk_rows(rows: list[dict[str, str]], columns: list[str]) -> list[list[dict[str, str]]]:
    if not rows:
        return [[]]

    col_widths = _column_widths(BODY_WIDTH_EMU)
    col_chars = _column_chars_per_line(col_widths)
    max_body_height = BODY_HEIGHT_EMU - HEADER_ROW_HEIGHT_EMU

    chunks: list[list[dict[str, str]]] = []
    current: list[dict[str, str]] = []
    current_height = 0

    for row in rows:
        row_height = _estimate_row_height(_row_values(row, columns, COLUMN_COUNT), col_chars)
        if current and current_height + row_height > max_body_height:
            chunks.append(current)
            current = []
            current_height = 0
        current.append(row)
        current_height += row_height

    if current:
        chunks.append(current)
    return chunks


def _delete_slide(prs: Presentation, index: int) -> None:
    slide_ids = prs.slides._sldIdLst
    slide_id = slide_ids[index]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    slide_ids.remove(slide_id)


def _move_slide_to_index(prs: Presentation, from_index: int, to_index: int) -> None:
    slide_ids = prs.slides._sldIdLst
    slide_id = slide_ids[from_index]
    slide_ids.remove(slide_id)
    slide_ids.insert(to_index, slide_id)


def _max_shape_id(slide) -> int:
    max_id = 1
    for shape in slide.shapes:
        shape_id = getattr(shape, "shape_id", None)
        if isinstance(shape_id, int):
            max_id = max(max_id, shape_id)
    return max_id


def _renumber_shape_ids(element, start_id: int) -> int:
    current = start_id
    for tag in _CNVPR_TAGS:
        for node in element.iter(tag):
            current += 1
            node.set("id", str(current))
    return current


def _duplicate_slide_safe(prs: Presentation, source_slide):
    dest = prs.slides.add_slide(source_slide.slide_layout)
    for shape in list(dest.shapes):
        dest.shapes._spTree.remove(shape._element)

    next_id = _max_shape_id(dest)
    for shape in source_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            continue
        new_element = deepcopy(shape.element)
        next_id = _renumber_shape_ids(new_element, next_id)
        dest.shapes._spTree.insert_element_before(new_element, "p:extLst")
    return dest


def _slide_title(slide) -> str | None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE:
            text = shape.text_frame.text.strip()
            return text.split("\n")[0].strip() if text else None
    return None


def _find_title_shape(slide):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE:
            return shape
    return None


def _find_main_body_shape(slide):
    bodies = [
        shape
        for shape in slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY
    ]
    if not bodies:
        return None
    return max(bodies, key=lambda shape: shape.width * shape.height)


def _oxml_child(parent, tag: str, **attrs) -> OxmlElement:
    element = OxmlElement(tag)
    element.attrib.update(attrs)
    parent.append(element)
    return element


def _strip_table_style(table) -> None:
    """Убираем встроенный tableStyle — иначе PowerPoint игнорирует границы ячеек."""
    tbl_pr = table._tbl.find(qn("a:tblPr"))
    if tbl_pr is None:
        return
    for node in list(tbl_pr.findall(qn("a:tableStyleId"))):
        tbl_pr.remove(node)
    for attr in ("firstRow", "bandRow", "lastRow", "firstCol", "lastCol", "bandCol"):
        tbl_pr.attrib.pop(attr, None)


def _set_cell_fill(cell, color_hex: str = "FFFFFF") -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for node in list(tc_pr.findall(qn("a:solidFill"))):
        tc_pr.remove(node)
    for node in list(tc_pr.findall(qn("a:noFill"))):
        tc_pr.remove(node)
    fill = _oxml_child(tc_pr, "a:solidFill")
    _oxml_child(fill, "a:srgbClr", val=color_hex.upper())


def _set_cell_outer_border(cell, color_hex: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        tag = qn(side)
        for existing in list(tc_pr.findall(tag)):
            tc_pr.remove(existing)
        line = _oxml_child(tc_pr, side, w=TABLE_BORDER_WIDTH, cap="flat", cmpd="sng", algn="ctr")
        fill = _oxml_child(line, "a:solidFill")
        _oxml_child(fill, "a:srgbClr", val=color_hex.upper())
        _oxml_child(line, "a:prstDash", val="solid")


def _set_cell_border(
    cell,
    *,
    color: str = TABLE_BORDER_COLOR,
    width: str = TABLE_BORDER_WIDTH,
) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        tag = qn(side)
        for existing in list(tc_pr.findall(tag)):
            tc_pr.remove(existing)
        line = _oxml_child(tc_pr, side, w=width, cap="flat", cmpd="sng", algn="ctr")
        fill = _oxml_child(line, "a:solidFill")
        _oxml_child(fill, "a:srgbClr", val=color)
        _oxml_child(line, "a:prstDash", val="solid")


def _apply_table_grid(table) -> None:
    for row in table.rows:
        for cell in row.cells:
            _set_cell_border(cell)


def _clear_secondary_body_placeholders(slide) -> None:
    bodies = [
        shape
        for shape in slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY
    ]
    if len(bodies) <= 1:
        return
    main = max(bodies, key=lambda shape: shape.width * shape.height)
    for shape in bodies:
        if shape is main:
            continue
        shape.text_frame.clear()


def _set_slide_title(shape, sheet_name: str) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = sheet_name
    run.font.name = TITLE_FONT_NAME
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def _fill_date_slide(slide, generated_at: datetime) -> None:
    formatted = generated_at.strftime("%d.%m.%Y")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if DATE_PLACEHOLDER not in text:
            continue
        shape.text_frame.text = text.replace(DATE_PLACEHOLDER, formatted)


def _find_body_with_markers(slide, *markers: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.is_placeholder:
            continue
        if shape.placeholder_format.type != PP_PLACEHOLDER.BODY:
            continue
        text = shape.text_frame.text
        if any(marker in text for marker in markers):
            return shape
    return None


def _capture_paragraph_font(paragraph) -> tuple[str | None, Pt | None, bool | None]:
    for run in paragraph.runs:
        if run.text.strip():
            return run.font.name, run.font.size, run.font.bold
    return None, None, None


def _capture_bullet_paragraph_pr(frame):
    for paragraph in frame.paragraphs:
        paragraph_pr = paragraph._p.find(qn("a:pPr"))
        if paragraph_pr is not None and paragraph_pr.find(qn("a:buChar")) is not None:
            return deepcopy(paragraph_pr)
    return None


def _apply_paragraph_pr(paragraph, paragraph_pr) -> None:
    existing = paragraph._p.find(qn("a:pPr"))
    if existing is not None:
        paragraph._p.remove(existing)
    paragraph._p.insert(0, deepcopy(paragraph_pr))


def _clear_paragraph_bullets(paragraph) -> None:
    paragraph_pr = paragraph._p.get_or_add_pPr()
    for tag in (
        "a:buChar",
        "a:buAutoNum",
        "a:buBlip",
        "a:buClr",
        "a:buClrTx",
        "a:buFont",
        "a:buSzPct",
        "a:buSzPts",
    ):
        for node in list(paragraph_pr.findall(qn(tag))):
            paragraph_pr.remove(node)
    if paragraph_pr.find(qn("a:buNone")) is None:
        paragraph_pr.insert(0, OxmlElement("a:buNone"))


def _append_notes_separator_paragraph(notes_frame, paragraph_count: int) -> int:
    paragraph = notes_frame.paragraphs[0] if paragraph_count == 0 else notes_frame.add_paragraph()
    _clear_paragraph_bullets(paragraph)
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_before = NOTES_BLOCK_SPACE_BEFORE
    paragraph.space_after = Pt(2)
    run = paragraph.add_run()
    run.text = NOTES_BLOCK_SEPARATOR
    run.font.name = TABLE_FONT_NAME
    run.font.size = TABLE_FONT_SIZE
    run.font.color.rgb = NOTES_MUTED_COLOR
    return paragraph_count + 1


def _fill_slide_notes(slide, rows: list[dict[str, str]], columns: list[str]) -> None:
    blocks = _slide_notes_blocks(rows, columns)
    if not blocks:
        return

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.word_wrap = True

    paragraph_count = 0
    for block_index, block in enumerate(blocks):
        if block_index > 0:
            paragraph_count = _append_notes_separator_paragraph(notes_frame, paragraph_count)

        for paragraph_index, text in enumerate(block):
            paragraph = notes_frame.paragraphs[0] if paragraph_count == 0 else notes_frame.add_paragraph()
            _clear_paragraph_bullets(paragraph)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_before = Pt(4) if paragraph_index > 0 else Pt(0)
            paragraph.space_after = (
                NOTES_PARAGRAPH_SPACE_AFTER if paragraph_index < len(block) - 1 else Pt(8)
            )
            run = paragraph.add_run()
            run.text = text
            run.font.name = TABLE_FONT_NAME
            run.font.size = TABLE_FONT_SIZE
            run.font.bold = paragraph_index == 0 and len(block) > 1
            run.font.color.rgb = TEXT_COLOR
            paragraph_count += 1


def _fill_body_lines(shape, lines: list[str], *, bullet: bool = False) -> None:
    frame = shape.text_frame
    font_name, font_size, font_bold = None, None, None
    bullet_pr = _capture_bullet_paragraph_pr(frame) if bullet else None
    for paragraph in frame.paragraphs:
        font_name, font_size, font_bold = _capture_paragraph_font(paragraph)
        if font_name or font_size:
            break

    frame.clear()
    frame.word_wrap = True
    if not lines:
        return

    for line_index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
        if bullet_pr is not None:
            _apply_paragraph_pr(paragraph, bullet_pr)
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = line
        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = font_size
        if font_bold is not None:
            run.font.bold = font_bold
        run.font.color.rgb = TEXT_COLOR


def _find_news_sheet(data: ProductStatusB2BOut) -> ProductStatusSheetOut | None:
    for sheet in data.sheets:
        if _normalize_title(sheet.name) == _normalize_title("Новости"):
            return sheet
    return data.sheets[0] if data.sheets else None


def _market_news_lines(sheet: ProductStatusSheetOut) -> list[str]:
    date_col = _find_column(sheet.columns, r"^Дата")
    news_col = _find_column(sheet.columns, r"^Новость")
    if not news_col:
        news_col = _find_column(sheet.columns, r"^Проект")
    desc_col = _find_column(sheet.columns, r"Описание")

    lines: list[str] = []
    for row in sheet.rows:
        date = display_cell_text(row.get(date_col, "") if date_col else "").strip()
        news = display_cell_text(row.get(news_col, "") if news_col else "").strip()
        description = display_cell_text(row.get(desc_col, "") if desc_col else "").strip()
        if not (date or news or description):
            continue
        lines.append(f"{date} | {news} | {description}")
    return lines


def _news_body_text_metrics(body_shape) -> tuple[int, int]:
    text_width = max(body_shape.width - NEWS_BODY_MARGIN_LEFT_EMU - CELL_TEXT_MARGIN_EMU, 1_000_000)
    chars_per_line = max(16, int(text_width / NEWS_CHAR_WIDTH_EMU))
    return chars_per_line, body_shape.height


def _estimate_news_item_height(line: str, chars_per_line: int) -> int:
    wrapped_lines = max(1, _estimate_text_lines(line, chars_per_line))
    return wrapped_lines * NEWS_LINE_HEIGHT_EMU + NEWS_PARAGRAPH_GAP_EMU


def _chunk_news_lines(lines: list[str], body_shape) -> list[list[str]]:
    if not lines:
        return [[]]

    chars_per_line, max_height = _news_body_text_metrics(body_shape)
    chunks: list[list[str]] = []
    current: list[str] = []
    current_height = 0

    for line in lines:
        item_height = _estimate_news_item_height(line, chars_per_line)
        if current and current_height + item_height > max_height:
            chunks.append(current)
            current = []
            current_height = 0
        current.append(line)
        current_height += item_height

    if current:
        chunks.append(current)
    return chunks


def _fill_market_news_slide_lines(slide, lines: list[str]) -> None:
    _clear_secondary_body_placeholders(slide)
    body_shape = _find_body_with_markers(slide, *MARKET_NEWS_MARKERS) or _find_main_body_shape(slide)
    if body_shape is None:
        return
    _fill_body_lines(body_shape, lines, bullet=True)


def _fill_market_news_slides(prs: Presentation, sheet: ProductStatusSheetOut) -> None:
    if len(prs.slides) <= MARKET_NEWS_SLIDE_INDEX:
        return

    lines = _market_news_lines(sheet)
    template_slide = prs.slides[MARKET_NEWS_SLIDE_INDEX]
    body_shape = _find_body_with_markers(template_slide, *MARKET_NEWS_MARKERS) or _find_main_body_shape(template_slide)
    if body_shape is None:
        return

    chunks = _chunk_news_lines(lines, body_shape)
    news_slides = [template_slide]
    insert_at = MARKET_NEWS_SLIDE_INDEX + 1

    for _ in chunks[1:]:
        duplicated_index = len(prs.slides)
        _duplicate_slide_safe(prs, template_slide)
        _move_slide_to_index(prs, duplicated_index, insert_at)
        news_slides.append(prs.slides[insert_at])
        insert_at += 1

    for slide, chunk in zip(news_slides, chunks, strict=True):
        _fill_market_news_slide_lines(slide, chunk)


def _try_load_market_news_sheet() -> ProductStatusSheetOut | None:
    try:
        news_data = load_b2b_news()
    except HTTPException:
        return None
    return _find_news_sheet(news_data)


def _hex_to_rgb(color_hex: str) -> RGBColor:
    value = color_hex.upper().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _set_run_strike(run, *, enabled: bool) -> None:
    r_pr = run._r.get_or_add_rPr()
    if enabled:
        r_pr.set("strike", "sngStrike")
    else:
        r_pr.attrib.pop("strike", None)


def _set_run_highlight(run, color_hex: str = DEFAULT_HIGHLIGHT_COLOR) -> None:
    r_pr = run._r.get_or_add_rPr()
    for child in list(r_pr):
        if child.tag == qn("a:highlight"):
            r_pr.remove(child)
    r_pr.insert(
        0,
        parse_xml(
            f'<a:highlight xmlns:a="{A_NS}">'
            f'<a:srgbClr val="{color_hex.upper()}"/></a:highlight>'
        ),
    )


def _apply_cell_container_style(cell, cell_style: CellStyle) -> None:
    if cell_style.bg:
        _set_cell_fill(cell, cell_style.bg)
    if cell_style.border:
        _set_cell_outer_border(cell, cell_style.border)


def _fill_white_cell(
    cell,
    value: str,
    *,
    col_index: int,
    bold: bool = False,
    default_text_color: RGBColor = TEXT_COLOR,
) -> None:
    cell_style, inner = split_cell_wrapper(value)
    sanitized = _sanitize_cell_text(inner)
    font_size, _ = _cell_font_size(col_index, sanitized)
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.margin_left = Pt(3)
    frame.margin_right = Pt(3)
    frame.margin_top = Pt(3)
    frame.margin_bottom = Pt(3)

    lines = sanitized.split("\n") if sanitized else [""]
    non_empty_lines = [line for line in lines if line.strip()]
    if not non_empty_lines:
        non_empty_lines = [""]

    for output_index, line in enumerate(non_empty_lines):
        paragraph = frame.paragraphs[0] if output_index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = CELL_PARAGRAPH_SPACE_BEFORE if output_index > 0 else Pt(0)
        paragraph.space_after = (
            CELL_PARAGRAPH_SPACE_AFTER if output_index < len(non_empty_lines) - 1 else Pt(0)
        )
        for segment in split_style_segments(line):
            if not segment.text:
                continue
            run = paragraph.add_run()
            run.text = segment.text
            run.font.name = TABLE_FONT_NAME
            run.font.size = font_size
            run.font.bold = bold or segment.bold
            run.font.italic = segment.italic
            run.font.color.rgb = _hex_to_rgb(segment.fg) if segment.fg else default_text_color
            _set_run_strike(run, enabled=segment.strike)
            if segment.bg and not segment.fg:
                _set_run_highlight(run, segment.bg)

    if not cell_style.bg:
        _set_cell_fill(cell, "FFFFFF")
    _apply_cell_container_style(cell, cell_style)


def _apply_table_column_widths(table, total_width: int) -> None:
    widths = _column_widths(total_width)
    for index, width in enumerate(widths):
        table.columns[index].width = width


def _layout_table_rows(table, rows: list[dict[str, str]], columns: list[str], total_height: int) -> None:
    col_widths = _column_widths(BODY_WIDTH_EMU)
    col_chars = _column_chars_per_line(col_widths)
    data_heights = [
        _estimate_row_height(_row_values(row, columns, COLUMN_COUNT), col_chars) for row in rows
    ]
    header_height = HEADER_ROW_HEIGHT_EMU
    available = max(total_height - header_height, MIN_ROW_HEIGHT_EMU * max(len(rows), 1))
    if not data_heights:
        table.rows[0].height = header_height
        return

    total_min = sum(data_heights)
    if total_min <= available:
        table.rows[0].height = header_height
        for index, height in enumerate(data_heights, start=1):
            table.rows[index].height = height
        return

    scale = available / total_min
    table.rows[0].height = header_height
    scaled = [max(MIN_ROW_HEIGHT_EMU, int(height * scale)) for height in data_heights]
    delta = available - sum(scaled)
    if delta:
        scaled[-1] += delta
    for index, height in enumerate(scaled, start=1):
        table.rows[index].height = height


def _fill_content_slide(
    slide,
    *,
    sheet_name: str,
    columns: list[str],
    rows: list[dict[str, str]],
) -> None:
    _clear_secondary_body_placeholders(slide)

    title_shape = _find_title_shape(slide)
    if title_shape is not None:
        _set_slide_title(title_shape, sheet_name)

    body_shape = _find_main_body_shape(slide)
    if body_shape is None:
        return

    mapped_headers = _presentation_headers(columns)
    data_rows = len(rows)
    table_shape = slide.shapes.add_table(
        data_rows + 1,
        COLUMN_COUNT,
        body_shape.left,
        body_shape.top,
        body_shape.width,
        body_shape.height,
    )
    table = table_shape.table
    _strip_table_style(table)
    _apply_table_column_widths(table, body_shape.width)

    for col_index, header in enumerate(mapped_headers):
        _fill_white_cell(table.rows[0].cells[col_index], header, col_index=col_index, bold=True)

    for row_index, row in enumerate(rows, start=1):
        values = _row_values(row, columns, COLUMN_COUNT)
        row_text_color = ATTENTION_TEXT_COLOR if _is_attention_row(row, columns) else TEXT_COLOR
        for col_index, value in enumerate(values):
            _fill_white_cell(
                table.rows[row_index].cells[col_index],
                value,
                col_index=col_index,
                default_text_color=row_text_color,
            )

    _apply_table_grid(table)
    _layout_table_rows(table, rows, columns, body_shape.height)

    _fill_slide_notes(slide, rows, columns)


def _build_slide_specs(
    data: ProductStatusB2BOut,
    catalog: TemplateCatalog,
) -> list[tuple[ProductStatusSheetOut, ContentSlideTemplate, list[dict[str, str]]]]:
    specs: list[tuple[ProductStatusSheetOut, ContentSlideTemplate, list[dict[str, str]]]] = []
    for sheet in data.sheets:
        rows = _filter_presentation_rows(sheet.rows, sheet.columns)
        if not rows:
            continue
        for template, chunk in catalog.chunk_plan(sheet.name, rows, sheet.columns):
            specs.append((sheet, template, chunk))
    return specs


def _read_presentation_sections(prs: Presentation) -> list[tuple[str, list[int]]]:
    return []


def generate_b2b_product_status_presentation(
    data: ProductStatusB2BOut | None = None,
) -> tuple[bytes, str]:
    if data is None:
        from app.db import SessionLocal, close_db_session

        db = SessionLocal()
        try:
            payload = load_b2b_product_status(db=db)
        finally:
            close_db_session(db)
    else:
        payload = data

    prs = _open_template_presentation()

    if len(prs.slides) < FIXED_SLIDE_COUNT + 1:
        raise HTTPException(
            status_code=503,
            detail="В шаблоне презентации должны быть титульный, служебные и контентные слайды.",
        )

    generated_at = datetime.now(MOSCOW_TZ)
    catalog = TemplateCatalog.from_presentation(prs)
    specs = _build_slide_specs(payload, catalog)
    if not specs:
        raise HTTPException(
            status_code=502,
            detail="Нет данных для формирования презентации.",
        )

    _fill_date_slide(prs.slides[0], generated_at)

    for sheet, template, chunk in specs:
        template_index = _template_index_for_sheet(sheet.name)
        source_slide = prs.slides[template_index]
        new_slide = _duplicate_slide_safe(prs, source_slide)
        _fill_content_slide(
            new_slide,
            sheet_name=sheet.name,
            columns=sheet.columns,
            rows=chunk,
        )

    for index in sorted(set(_SHEET_TEMPLATE_INDEX.values()), reverse=True):
        if index >= FIXED_SLIDE_COUNT:
            _delete_slide(prs, index)

    news_sheet = _try_load_market_news_sheet()
    if news_sheet is not None:
        _fill_market_news_slides(prs, news_sheet)

    buffer = io.BytesIO()
    prs.save(buffer)
    filename = f"status-produkta-b2b-{generated_at.strftime('%Y%m%d')}.pptx"
    return buffer.getvalue(), filename
