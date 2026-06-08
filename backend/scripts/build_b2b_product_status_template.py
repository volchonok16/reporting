#!/usr/bin/env python3
"""Собирает чистый шаблон PPTX для выгрузки статуса продукта B2B."""

from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_SOURCE = ROOT / "assets" / "b2b_product_status_template.pptx"
DEFAULT_OUTPUT = ROOT / "assets" / "b2b_product_status_template.pptx"
COVER_SLIDE_INDEX = 0
MIN_ROW_HEIGHT_EMU = 396000


def _remove_shape(shape) -> None:
    element = shape.element
    element.getparent().remove(element)


def _fix_table_row_heights(slide) -> int:
    fixed = 0
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        row_heights = [int(row.height) for row in table.rows]
        if not row_heights:
            continue
        positive = [height for height in row_heights if height > 0]
        fallback = max(positive) if positive else MIN_ROW_HEIGHT_EMU
        for row in table.rows:
            if int(row.height) <= 0:
                row.height = fallback
                fixed += 1
    return fixed


def _clean_cover_slide(slide) -> None:
    title_shapes: list = []
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            _remove_shape(shape)
            continue
        if shape.has_table:
            _remove_shape(shape)
            continue
        if shape.has_text_frame and "заголовок" in shape.name.lower():
            title_shapes.append(shape)

    for extra_title in title_shapes[1:]:
        _remove_shape(extra_title)


def build_template(*, source: Path, output: Path) -> None:
    if not source.is_file():
        raise FileNotFoundError(f"Исходный шаблон не найден: {source}")

    if source.resolve() != output.resolve():
        output.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, output)

    prs = Presentation(str(output))
    removed_ole = 0
    fixed_rows = 0

    for index, slide in enumerate(prs.slides):
        if index == COVER_SLIDE_INDEX:
            _clean_cover_slide(slide)

        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
                _remove_shape(shape)
                removed_ole += 1

        fixed_rows += _fix_table_row_heights(slide)

    prs.save(str(output))
    print(f"saved={output}")
    print(f"slides={len(prs.slides)} removed_ole={removed_ole} fixed_rows={fixed_rows}")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source", type=Path, default=DEFAULT_SOURCE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args(argv)
    build_template(source=args.source, output=args.output)
    return 0


if __name__ == "__main__":
    sys.exit(main())
