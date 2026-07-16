import io
import re
from types import SimpleNamespace
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Pt

from app.product_status_presentation import (
    DATE_PLACEHOLDER,
    DESCRIPTION_HEADER,
    FIXED_SLIDE_COUNT,
    MARKET_NEWS_SLIDE_INDEX,
    NOTES_BLOCK_SEPARATOR,
    TABLE_FONT_SIZE,
    TABLE_FONT_SIZE_DENSE,
    TemplateCatalog,
    _cell_font_size,
    _chunk_news_lines,
    _chunk_rows,
    _description_value_column,
    _why_value_column,
    _estimate_row_height,
    _estimate_text_lines,
    _attention_column,
    _is_attention_row,
    _filter_presentation_rows,
    _mapped_columns,
    _market_news_lines,
    _normalize_title,
    _presentation_headers,
    _row_values,
    _section_name_for_sheet,
    _slide_notes_blocks,
    _slide_notes_text,
    _template_index_for_sheet,
    generate_b2b_product_status_presentation,
)
from app.product_status_presentation import _column_chars_per_line, _column_widths, BODY_WIDTH_EMU
from app.schemas import ProductStatusB2BOut, ProductStatusSheetOut

TEMPLATE_PATH = (
    __import__("pathlib").Path(__file__).resolve().parents[1]
    / "assets"
    / "b2b_product_status_template.pptx"
)


def test_template_has_status_structure() -> None:
    prs = Presentation(str(TEMPLATE_PATH))
    assert len(prs.slides) >= 8
    texts = [shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame]
    assert DATE_PLACEHOLDER in "".join(texts)


def test_cell_font_size_matches_template_columns() -> None:
    assert _cell_font_size(0, "июль") == (TABLE_FONT_SIZE, "1000")
    assert _cell_font_size(3, "Кратко") == (TABLE_FONT_SIZE_DENSE, "800")


def test_normalize_title() -> None:
    assert _normalize_title("  Продуктовый офис: CORE  ") == "продуктовый офис: core"


def test_section_name_for_sheet() -> None:
    assert _section_name_for_sheet("Продуктовый офис: SMS") == "SMS"
    assert _section_name_for_sheet("Офис: SMS") == "SMS"
    assert _section_name_for_sheet("Продуктовый офис: VOICE") == "Voice"
    assert _section_name_for_sheet("Офис: VOICE") == "Voice"


def test_template_index_for_sheet_colors() -> None:
    assert _template_index_for_sheet("Продуктовый офис: CORE") == 3
    assert _template_index_for_sheet("Офис: CORE") == 3
    assert _template_index_for_sheet("Продуктовый офис: M2M / IoT") == 4
    assert _template_index_for_sheet("Офис: M2M / IoT") == 4
    assert _template_index_for_sheet("Продуктовый офис: SMS") == 5
    assert _template_index_for_sheet("Офис: SMS") == 5
    assert _template_index_for_sheet("Продуктовый офис: VOICE") == 6
    assert _template_index_for_sheet("Офис: VOICE") == 6
    assert _template_index_for_sheet("Продуктовый офис: Перспективные продукты") == 7
    assert _template_index_for_sheet("Офис: Перспективные продукты") == 7
    assert _template_index_for_sheet("Продуктовый офис: Продуктовый маркетинг") == 3
    assert _template_index_for_sheet("Офис: Продуктовый маркетинг") == 3


def test_mapped_columns_falls_back_to_column_order() -> None:
    columns = ["Месяц", "Направление", "Статус работ", "Комментарий"]
    assert _mapped_columns(columns) == columns


def test_row_values_maps_semantic_columns() -> None:
    columns = [
        "Дата запуска",
        "Проект",
        "Описание проекта и статус",
        "Зачем и для чего делаем",
    ]
    row = {
        "Дата запуска": "01.06",
        "Проект": "SMS Hub",
        "Описание проекта и статус": "Статус",
        "Зачем и для чего делаем": "Комментарий",
    }
    assert _row_values(row, columns, 4) == ["01.06", "SMS Hub", "Статус", "Комментарий"]


def test_row_values_uses_presentation_description_column() -> None:
    columns = [
        "ЗНИ",
        "Идет в презентацию",
        "Обратить внимание",
        "Дата запуска",
        "Проект",
        "Полное Описание проекта и статус",
        "Для презентации Описание проекта и статус",
        "Описание проекта и статус",
        "Зачем и для чего делаем полное описание",
        "Зачем и для чего делаем для презентации",
        "Зачем и для чего делаем",
    ]
    row = {
        "Дата запуска": "01.06",
        "Проект": "SMS Hub",
        "Полное Описание проекта и статус": "Полный текст",
        "Для презентации Описание проекта и статус": "Коротко для слайда",
        "Описание проекта и статус": "Не использовать",
        "Зачем и для чего делаем полное описание": "Полный зачем",
        "Зачем и для чего делаем для презентации": "Короткий зачем",
        "Зачем и для чего делаем": "Не использовать",
    }
    assert _description_value_column(columns) == "Для презентации Описание проекта и статус"
    assert _why_value_column(columns) == "Зачем и для чего делаем для презентации"
    assert _row_values(row, columns, 4) == [
        "01.06",
        "SMS Hub",
        "Коротко для слайда",
        "Короткий зачем",
    ]
    assert _presentation_headers(columns)[2] == DESCRIPTION_HEADER


def test_filter_presentation_rows_keeps_only_yes() -> None:
    columns = ["Идет в презентацию", "Проект"]
    rows = [
        {"Идет в презентацию": "Да", "Проект": "A"},
        {"Идет в презентацию": "Нет", "Проект": "B"},
        {"Идет в презентацию": "", "Проект": "C"},
        {"Идет в презентацию": "да", "Проект": "D"},
        {"Идет в презентацию": "<<cell:bg:C6EFCE>>да<<>>", "Проект": "E"},
    ]
    assert _filter_presentation_rows(rows, columns) == [
        {"Идет в презентацию": "Да", "Проект": "A"},
        {"Идет в презентацию": "да", "Проект": "D"},
        {"Идет в презентацию": "<<cell:bg:C6EFCE>>да<<>>", "Проект": "E"},
    ]


def test_attention_column_and_row() -> None:
    columns = ["Идет в презентацию", "Обратить внимание", "Проект"]
    assert _attention_column(columns) == "Обратить внимание"
    assert _is_attention_row({"Обратить внимание": "Да", "Проект": "A"}, columns)
    assert not _is_attention_row({"Обратить внимание": "Нет", "Проект": "B"}, columns)
    assert not _is_attention_row({"Проект": "C"}, ["Проект"])


def test_slide_notes_text_joins_full_descriptions() -> None:
    columns = ["Проект", "Полное Описание проекта и статус"]
    rows = [
        {"Проект": "CORE", "Полное Описание проекта и статус": "Первая заметка"},
        {"Проект": "SMS", "Полное Описание проекта и статус": "Вторая заметка"},
    ]
    assert _slide_notes_text(rows, columns) == (
        "CORE\n\nПервая заметка\n\n—\n\nSMS\n\nВторая заметка"
    )


def test_slide_notes_blocks_splits_multiline_descriptions() -> None:
    columns = ["Проект", "Полное Описание проекта и статус"]
    rows = [
        {
            "Проект": "SMS Hub",
            "Полное Описание проекта и статус": "Первый абзац\n\nВторой абзац",
        },
    ]
    assert _slide_notes_blocks(rows, columns) == [
        ["SMS Hub", "Первый абзац\n\nВторой абзац"],
    ]


def test_slide_notes_and_table_use_distinct_column_values() -> None:
    columns = [
        "Дата запуска",
        "Проект",
        "Полное Описание проекта и статус",
        "Для презентации Описание проекта и статус",
        "Описание проекта и статус",
        "Зачем и для чего делаем полное описание",
        "Зачем и для чего делаем для презентации",
        "Зачем и для чего делаем",
    ]
    row = {
        "Дата запуска": "01.06",
        "Проект": "SMS Hub",
        "Полное Описание проекта и статус": "Полный статус",
        "Для презентации Описание проекта и статус": "Короткий статус",
        "Описание проекта и статус": "Легаси",
        "Зачем и для чего делаем полное описание": "Полный зачем",
        "Зачем и для чего делаем для презентации": "Короткий зачем",
        "Зачем и для чего делаем": "Легаси зачем",
    }
    assert _row_values(row, columns, 4) == [
        "01.06",
        "SMS Hub",
        "Короткий статус",
        "Короткий зачем",
    ]
    assert _slide_notes_blocks([row], columns) == [
        ["SMS Hub", "Полный статус", "Полный зачем"],
    ]


def test_slide_notes_blocks_include_full_why_text() -> None:
    columns = [
        "Проект",
        "Полное Описание проекта и статус",
        "Зачем и для чего делаем полное описание",
    ]
    rows = [
        {
            "Проект": "CORE",
            "Полное Описание проекта и статус": "Полный статус",
            "Зачем и для чего делаем полное описание": "Полный зачем",
        },
    ]
    assert _slide_notes_blocks(rows, columns) == [
        ["CORE", "Полный статус", "Полный зачем"],
    ]


def test_slide_notes_blocks_support_why_only_rows() -> None:
    columns = ["Проект", "Зачем и для чего делаем полное описание"]
    rows = [
        {
            "Проект": "SMS",
            "Зачем и для чего делаем полное описание": "Только зачем в заметках",
        },
    ]
    assert _slide_notes_text(rows, columns) == "SMS\n\nТолько зачем в заметках"


def test_slide_notes_blocks_one_table_row_per_block() -> None:
    columns = ["Проект", "Полное Описание проекта и статус"]
    rows = [
        {
            "Проект": "SMS Hub",
            "Полное Описание проекта и статус": "Тендер ОТП-банк\nКейс с Озоном",
        },
        {
            "Проект": "SMS Hub",
            "Полное Описание проекта и статус": "Кастомные отчеты",
        },
    ]
    assert _slide_notes_blocks(rows, columns) == [
        ["SMS Hub", "Тендер ОТП-банк\nКейс с Озоном"],
        ["SMS Hub", "Кастомные отчеты"],
    ]


def test_slide_notes_blocks_use_legacy_columns_without_presentation_split() -> None:
    columns = ["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"]
    rows = [
        {
            "Дата запуска": "01.06",
            "Проект": "CORE",
            "Описание проекта": "Полный статус для заметок",
            "Зачем и для чего делаем": "Полный зачем для заметок",
        },
    ]
    assert _slide_notes_blocks(rows, columns) == [
        ["CORE", "Полный статус для заметок", "Полный зачем для заметок"],
    ]


def test_slide_notes_blocks_fallback_to_legacy_when_full_cells_empty() -> None:
    columns = [
        "Проект",
        "Полное Описание проекта и статус",
        "Для презентации Описание проекта и статус",
        "Описание проекта и статус",
        "Зачем и для чего делаем полное описание",
        "Зачем и для чего делаем для презентации",
        "Зачем и для чего делаем",
    ]
    rows = [
        {
            "Проект": "SMS Hub",
            "Полное Описание проекта и статус": "",
            "Для презентации Описание проекта и статус": "Коротко",
            "Описание проекта и статус": "Полный текст из соседнего столбца",
            "Зачем и для чего делаем полное описание": "",
            "Зачем и для чего делаем для презентации": "Короткий зачем",
            "Зачем и для чего делаем": "Полный зачем из соседнего столбца",
        },
    ]
    assert _slide_notes_blocks(rows, columns) == [
        [
            "SMS Hub",
            "Полный текст из соседнего столбца",
            "Полный зачем из соседнего столбца",
        ],
    ]


def test_slide_notes_blocks_use_legacy_why_without_full_description_column() -> None:
    columns = [
        "Проект",
        "Для презентации Описание проекта и статус",
        "Зачем и для чего делаем для презентации",
        "Зачем и для чего делаем",
    ]
    rows = [
        {
            "Проект": "CORE",
            "Для презентации Описание проекта и статус": "Коротко",
            "Зачем и для чего делаем для презентации": "Короткий зачем",
            "Зачем и для чего делаем": "Полный зачем в заметках",
        },
    ]
    assert _slide_notes_text(rows, columns) == "CORE\n\nПолный зачем в заметках"


def test_slide_notes_blocks_skip_unified_why_when_description_split_exists() -> None:
    columns = [
        "Проект",
        "Полное Описание проекта и статус",
        "Для презентации Описание проекта и статус",
        "Зачем и для чего делаем",
    ]
    row = {
        "Проект": "SMS Hub",
        "Полное Описание проекта и статус": "Полный статус",
        "Для презентации Описание проекта и статус": "Короткий статус",
        "Зачем и для чего делаем": "Зачем на слайде",
    }
    assert _why_value_column(columns) == "Зачем и для чего делаем"
    assert _row_values(row, columns, 4)[3] == "Зачем на слайде"
    assert _slide_notes_blocks([row], columns) == [
        ["SMS Hub", "Полный статус"],
    ]


def test_estimate_text_lines_wraps_each_paragraph_separately() -> None:
    short_lines = "Первая строка\nВторая строка\nТретья строка"
    assert _estimate_text_lines(short_lines, 40) == 3

    long_paragraph = " ".join(["Фраза с пояснением"] * 18)
    col_chars = _column_chars_per_line(_column_widths(BODY_WIDTH_EMU))
    dense_chars = col_chars[3]
    per_line_wrapping = max(1, (len(long_paragraph) + dense_chars - 1) // dense_chars)
    assert _estimate_text_lines(long_paragraph, dense_chars) == per_line_wrapping


def test_estimate_row_height_uses_dense_column_metrics() -> None:
    columns = ["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"]
    row = {
        "Дата запуска": "июль",
        "Проект": "KATC",
        "Описание проекта": "Короткий статус",
        "Зачем и для чего делаем": (
            "Первая строка пояснения. "
            "Вторая строка пояснения. "
            "Третья строка пояснения. "
            "ориентировочный срок тестирования 15 июня"
        ),
    }
    col_chars = _column_chars_per_line(_column_widths(BODY_WIDTH_EMU))
    height = _estimate_row_height(_row_values(row, columns, 4), col_chars)
    lines = _estimate_text_lines(row["Зачем и для чего делаем"], col_chars[3])
    assert lines <= 4
    assert height < lines * 145000 + 45000


def test_chunk_rows_splits_long_sheet() -> None:
    columns = ["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"]
    rows = [
        {
            "Дата запуска": "01.06",
            "Проект": "CORE",
            "Описание проекта": "Коротко",
            "Зачем и для чего делаем": "OK",
        },
        {
            "Дата запуска": "02.06",
            "Проект": "CORE",
            "Описание проекта": "Очень длинное описание " * 120,
            "Зачем и для чего делаем": "Ещё текст " * 80,
        },
    ]
    chunks = _chunk_rows(rows, columns)
    assert sum(len(chunk) for chunk in chunks) == 2
    assert len(chunks) >= 2


def test_catalog_chunk_plan_keeps_sheet_order() -> None:
    prs = Presentation(str(TEMPLATE_PATH))
    catalog = TemplateCatalog.from_presentation(prs)
    rows = [{"Дата запуска": "01.06", "Проект": "CORE", "Описание проекта": "A", "Зачем и для чего делаем": "B"}]
    chunks = catalog.chunk_plan("Продуктовый офис: CORE", rows, ["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"])
    assert chunks[0][0].index == 3
    assert len(chunks[0][1]) == 1


def test_market_news_lines_formats_rows() -> None:
    sheet = ProductStatusSheetOut(
        gid="0",
        name="Новости",
        columns=["Дата", "Новость", "Описание"],
        rows=[
            {"Дата": "08.06", "Новость": "МТС изменил тариф", "Описание": "Подробности"},
            {"Дата": "", "Новость": "", "Описание": ""},
            {"Дата": "09.06", "Новость": "Yota обновила планы", "Описание": "Стоимость выросла"},
        ],
        totalShown=3,
    )
    assert _market_news_lines(sheet) == [
        "08.06 | МТС изменил тариф | Подробности",
        "09.06 | Yota обновила планы | Стоимость выросла",
    ]


def test_chunk_news_lines_splits_when_body_is_short() -> None:
    body = SimpleNamespace(width=4_474_200, height=900_000)
    lines = [f"01.06 | Новость {index} | " + "подробность " * 120 for index in range(6)]
    chunks = _chunk_news_lines(lines, body)
    assert len(chunks) >= 2
    assert sum(len(chunk) for chunk in chunks) == 6


@patch("app.product_status_presentation.load_b2b_news")
@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_splits_market_news_across_slides(mock_load, mock_news) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[
                    {
                        "Дата запуска": "",
                        "Проект": "CORE",
                        "Описание проекта": "Коротко",
                        "Зачем и для чего делаем": "OK",
                    }
                ],
                totalShown=1,
            )
        ],
    )
    mock_news.return_value = ProductStatusB2BOut(
        title="Новости и запуски",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Новости",
                columns=["Дата", "Новость", "Описание"],
                rows=[
                    {
                        "Дата": f"0{index}.06",
                        "Новость": f"Новость {index}",
                        "Описание": "подробность " * 300,
                    }
                    for index in range(1, 5)
                ],
                totalShown=4,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))

    news_slide_count = 0
    for index in range(MARKET_NEWS_SLIDE_INDEX, len(prs.slides)):
        slide = prs.slides[index]
        if any(shape.has_table for shape in slide.shapes):
            break
        joined = "".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        )
        if "<$start>" in joined:
            break
        if " | " in joined and "<$date>" not in joined:
            news_slide_count += 1
        else:
            break

    assert news_slide_count >= 2
    assert len(prs.slides) > FIXED_SLIDE_COUNT + 1


@patch("app.product_status_presentation.load_b2b_news")
@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_fills_market_news_slide(mock_load, mock_news) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[
                    {
                        "Дата запуска": "",
                        "Проект": "CORE",
                        "Описание проекта": "Перенос номеров",
                        "Зачем и для чего делаем": "В работе",
                    }
                ],
                totalShown=1,
            )
        ],
    )
    mock_news.return_value = ProductStatusB2BOut(
        title="Новости и запуски",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Новости",
                columns=["Дата", "Новость", "Описание"],
                rows=[
                    {
                        "Дата": "08.06",
                        "Новость": "МТС изменил тариф",
                        "Описание": "Подробности",
                    },
                    {
                        "Дата": "09.06",
                        "Новость": "Yota обновила планы",
                        "Описание": "Стоимость выросла",
                    },
                ],
                totalShown=2,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    slide = prs.slides[MARKET_NEWS_SLIDE_INDEX]
    body_texts = [
        shape.text_frame.text
        for shape in slide.shapes
        if shape.has_text_frame and shape.is_placeholder and shape.text_frame.text.strip()
    ]
    assert len(body_texts) == 1
    assert "<$date>" not in body_texts[0]
    assert "08.06 | МТС изменил тариф | Подробности" in body_texts[0]
    assert "09.06 | Yota обновила планы | Стоимость выросла" in body_texts[0]

    body_shape = next(
        shape
        for shape in slide.shapes
        if shape.has_text_frame and shape.is_placeholder and shape.text_frame.text.strip()
    )
    paragraphs = [paragraph.text for paragraph in body_shape.text_frame.paragraphs if paragraph.text.strip()]
    assert len(paragraphs) == 2
    assert paragraphs[0] == "08.06 | МТС изменил тариф | Подробности"
    assert paragraphs[1] == "09.06 | Yota обновила планы | Стоимость выросла"

    slide_xml = slide.part.blob.decode()
    first_paragraph_start = slide_xml.index("<a:t>08.06 | МТС изменил тариф | Подробности</a:t>")
    first_paragraph_block = slide_xml[max(0, first_paragraph_start - 800):first_paragraph_start]
    assert 'char="•"' in first_paragraph_block or "<a:buChar" in first_paragraph_block


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_fills_date_and_keeps_fixed_slides(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[
                    {
                        "Дата запуска": "",
                        "Проект": "CORE",
                        "Описание проекта": "Перенос номеров",
                        "Зачем и для чего делаем": "В работе",
                    }
                ],
                totalShown=1,
            )
        ],
    )

    content, filename = generate_b2b_product_status_presentation()
    assert filename.startswith("status-produkta-b2b-")
    assert filename.endswith(".pptx")

    prs = Presentation(io.BytesIO(content))
    assert len(prs.slides) == FIXED_SLIDE_COUNT + 1

    slide1_text = "".join(shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame)
    assert "Статус B2B" in slide1_text
    assert DATE_PLACEHOLDER not in slide1_text
    assert re.search(r"\d{2}\.\d{2}\.\d{4}", slide1_text)

    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    title = next(
        shape.text_frame.text
        for shape in content_slide.shapes
        if shape.has_text_frame and shape.is_placeholder
    )
    assert title == "Продуктовый офис: CORE"

    table = next(shape.table for shape in content_slide.shapes if shape.has_table)
    assert table.rows[1].cells[1].text == "CORE"
    assert table.rows[0].cells[0].text == "Дата запуска"


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_uses_colored_templates_per_sheet(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[{"Дата запуска": "", "Проект": "CORE", "Описание проекта": "A", "Зачем и для чего делаем": "B"}],
                totalShown=1,
            ),
            ProductStatusSheetOut(
                gid="1",
                name="Продуктовый офис: SMS",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[{"Дата запуска": "", "Проект": "SMS Hub", "Описание проекта": "C", "Зачем и для чего делаем": "D"}],
                totalShown=1,
            ),
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))

    core_slide = prs.slides[FIXED_SLIDE_COUNT]
    sms_slide = next(
        prs.slides[index]
        for index in range(FIXED_SLIDE_COUNT, len(prs.slides))
        if any(
            shape.has_text_frame and shape.text_frame.text == "Продуктовый офис: SMS"
            for shape in prs.slides[index].shapes
        )
    )
    core_layout = core_slide.slide_layout.name
    sms_layout = sms_slide.slide_layout.name
    assert core_layout
    assert sms_layout
    assert core_slide.slide_layout == prs.slides[FIXED_SLIDE_COUNT].slide_layout
    assert sms_slide.slide_layout != core_slide.slide_layout


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_colored_text_has_no_highlight_fill(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: SMS",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[
                    {
                        "Дата запуска": "",
                        "Проект": "[[fg:FF0000::Контроль]]",
                        "Описание проекта": "[[fg:1254CC::hh.ru]] — риск оттока",
                        "Зачем и для чего делаем": "",
                    }
                ],
                totalShown=1,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    xml = content_slide.part.blob.decode()
    control_start = xml.index("<a:t>Контроль</a:t>")
    control_run = xml[max(0, control_start - 500):control_start + 200]
    assert "val=\"FF0000\"" in control_run
    assert "<a:highlight>" not in control_run

    link_start = xml.index("<a:t>hh.ru</a:t>")
    link_run = xml[max(0, link_start - 500):link_start + 200]
    assert "val=\"1254CC\"" in link_run
    assert "<a:highlight>" not in link_run


@patch("app.product_status_presentation._try_load_market_news_sheet", return_value=None)
@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_preserves_toolbar_styles(mock_load, _mock_news) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=[
                    "Идет в презентацию",
                    "Дата запуска",
                    "Проект",
                    "Для презентации Описание проекта и статус",
                    "Зачем и для чего делаем для презентации",
                ],
                rows=[
                    {
                        "Идет в презентацию": "Да",
                        "Дата запуска": "01.06",
                        "Проект": "[[bold::Жирный]] [[italic::Курсив]]",
                        "Для презентации Описание проекта и статус": (
                            "$Жёлтый$ {{C6EFCE:Зелёный}} {{FFC7CE:Розовый}} {{BDD7EE:Голубой}} "
                            "[[fg:B91C1C::Красный]] [[fg:0000FF::Синий]] [[fg:000000::Чёрный]] "
                            "[[fg:008000::Зелёный текст]] [[fg:808080::Серый]] "
                            "[[fg:0000FF;strike::Зачёркнутый]]"
                        ),
                        "Зачем и для чего делаем для презентации": (
                            "[[bg:FFFF00;fg:0000FF::Маркер и цвет]]"
                        ),
                    }
                ],
                totalShown=1,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    xml = content_slide.part.blob.decode()

    for label in (
        "Жирный",
        "Курсив",
        "Жёлтый",
        "Зелёный",
        "Розовый",
        "Голубой",
        "Красный",
        "Синий",
        "Чёрный",
        "Зелёный текст",
        "Серый",
        "Зачёркнутый",
        "Маркер и цвет",
    ):
        assert f"<a:t>{label}</a:t>" in xml

    bold_start = xml.index("<a:t>Жирный</a:t>")
    bold_run = xml[max(0, bold_start - 500):bold_start + 200]
    assert 'b="1"' in bold_run

    italic_start = xml.index("<a:t>Курсив</a:t>")
    italic_run = xml[max(0, italic_start - 500):italic_start + 200]
    assert 'i="1"' in italic_run

    strike_start = xml.index("<a:t>Зачёркнутый</a:t>")
    strike_run = xml[max(0, strike_start - 500):strike_start + 200]
    assert 'strike="sngStrike"' in strike_run

    combo_start = xml.index("<a:t>Маркер и цвет</a:t>")
    combo_run = xml[max(0, combo_start - 500):combo_start + 200]
    assert 'val="FFFF00"' in combo_run
    assert 'val="0000FF"' in combo_run
    assert "<a:highlight>" in combo_run


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_attention_rows_use_red_text(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: SMS",
                columns=[
                    "Идет в презентацию",
                    "Обратить внимание",
                    "Дата запуска",
                    "Проект",
                    "Описание проекта",
                    "Зачем и для чего делаем",
                ],
                rows=[
                    {
                        "Идет в презентацию": "Да",
                        "Обратить внимание": "Да",
                        "Дата запуска": "09.06",
                        "Проект": "Критичный проект",
                        "Описание проекта": "Срочный статус",
                        "Зачем и для чего делаем": "Риск",
                    },
                    {
                        "Идет в презентацию": "Да",
                        "Обратить внимание": "Нет",
                        "Дата запуска": "10.06",
                        "Проект": "Обычный проект",
                        "Описание проекта": "Штатный статус",
                        "Зачем и для чего делаем": "OK",
                    },
                ],
                totalShown=2,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    xml = content_slide.part.blob.decode()

    critical_start = xml.index("<a:t>Критичный проект</a:t>")
    critical_run = xml[max(0, critical_start - 500):critical_start + 200]
    assert 'val="C00000"' in critical_run

    normal_start = xml.index("<a:t>Обычный проект</a:t>")
    normal_run = xml[max(0, normal_start - 500):normal_start + 200]
    assert 'val="C00000"' not in normal_run
    assert 'val="000000"' in normal_run


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_table_cells_contain_text_xml(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[
                    {
                        "Дата запуска": "09.06",
                        "Проект": "Ремонт",
                        "Описание проекта": "Убираем $300 рублевые офферы$",
                        "Зачем и для чего делаем": "В работе",
                    }
                ],
                totalShown=1,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    title = next(
        shape.text_frame.text
        for shape in content_slide.shapes
        if shape.is_placeholder and shape.has_text_frame
    )
    table = next(shape.table for shape in content_slide.shapes if shape.has_table)
    assert table.rows[1].cells[1].text == "Ремонт"
    assert table.rows[0].cells[0].text == "Дата запуска"
    xml = content_slide.part.blob.decode()

    assert "<a:t>Ремонт</a:t>" in xml
    assert "<a:t>300 рублевые офферы</a:t>" in xml
    assert "<a:highlight>" in xml or "strike=" in xml
    assert 'srgbClr val="FFFFFF"' in xml
    assert 'srgbClr val="7F7F7F"' in xml
    assert "tableStyleId" not in xml.split("<a:tbl>")[1].split("</a:tbl>")[0]
    assert title == "Продуктовый офис: CORE"


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_filters_rows_and_fills_notes(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=[
                    "Идет в презентацию",
                    "Дата запуска",
                    "Проект",
                    "Полное Описание проекта и статус",
                    "Для презентации Описание проекта и статус",
                    "Описание проекта и статус",
                    "Зачем и для чего делаем полное описание",
                    "Зачем и для чего делаем для презентации",
                    "Зачем и для чего делаем",
                ],
                rows=[
                    {
                        "Идет в презентацию": "Да",
                        "Дата запуска": "01.06",
                        "Проект": "CORE",
                        "Полное Описание проекта и статус": "Полный текст для заметок",
                        "Для презентации Описание проекта и статус": "Коротко",
                        "Описание проекта и статус": "Игнор",
                        "Зачем и для чего делаем полное описание": "Полный зачем для заметок",
                        "Зачем и для чего делаем для презентации": "Короткий зачем",
                        "Зачем и для чего делаем": "OK",
                    },
                    {
                        "Идет в презентацию": "Нет",
                        "Дата запуска": "02.06",
                        "Проект": "SKIP",
                        "Полное Описание проекта и статус": "Не попасть",
                        "Для презентации Описание проекта и статус": "Не попасть",
                        "Описание проекта и статус": "Не попасть",
                        "Зачем и для чего делаем": "SKIP",
                    },
                ],
                totalShown=2,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    table = next(shape.table for shape in content_slide.shapes if shape.has_table)
    assert table.rows[0].cells[2].text == DESCRIPTION_HEADER
    assert table.rows[1].cells[1].text == "CORE"
    assert table.rows[1].cells[2].text == "Коротко"
    assert table.rows[1].cells[3].text == "Короткий зачем"
    assert len(table.rows) == 2
    assert "SKIP" not in content_slide.part.blob.decode()
    assert "Игнор" not in content_slide.part.blob.decode()
    assert "OK" not in table.rows[1].cells[3].text

    notes = content_slide.notes_slide.notes_text_frame.text
    assert "CORE" in notes
    assert "Полный текст для заметок" in notes
    assert "Полный зачем для заметок" in notes
    assert "Коротко" not in notes
    assert "Короткий зачем" not in notes
    assert "Не попасть" not in notes
    assert NOTES_BLOCK_SEPARATOR not in notes

    notes_xml = content_slide.notes_slide.part.blob.decode()
    assert "<a:buChar" not in notes_xml
    assert "<a:buNone" in notes_xml


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_notes_use_block_separators(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="1512199647",
                name="Продуктовый офис: SMS",
                columns=[
                    "Идет в презентацию",
                    "Дата запуска",
                    "Проект",
                    "Полное Описание проекта и статус",
                    "Для презентации Описание проекта и статус",
                    "Зачем и для чего делаем полное описание",
                    "Зачем и для чего делаем для презентации",
                ],
                rows=[
                    {
                        "Идет в презентацию": "Да",
                        "Дата запуска": "01.06",
                        "Проект": "SMS Hub",
                        "Полное Описание проекта и статус": "Тендер ОТП-банк",
                        "Для презентации Описание проекта и статус": "Тендер",
                        "Зачем и для чего делаем полное описание": "Полный зачем 1",
                        "Зачем и для чего делаем для презентации": "Короткий зачем 1",
                    },
                    {
                        "Идет в презентацию": "Да",
                        "Дата запуска": "02.06",
                        "Проект": "SMS Hub",
                        "Полное Описание проекта и статус": "Кейс с Озоном",
                        "Для презентации Описание проекта и статус": "Озон",
                        "Зачем и для чего делаем полное описание": "Полный зачем 2",
                        "Зачем и для чего делаем для презентации": "Короткий зачем 2",
                    },
                ],
                totalShown=2,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    notes = content_slide.notes_slide.notes_text_frame.text
    assert "—" in notes
    assert "Тендер ОТП-банк" in notes
    assert "Полный зачем 1" in notes
    assert "Кейс с Озоном" in notes
    assert "Полный зачем 2" in notes
    assert "Короткий зачем 1" not in notes
    assert "Короткий зачем 2" not in notes


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_splits_rows_across_slides(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[
                    {
                        "Дата запуска": "01.06",
                        "Проект": "CORE",
                        "Описание проекта": "Коротко",
                        "Зачем и для чего делаем": "OK",
                    },
                    {
                        "Дата запуска": "02.06",
                        "Проект": "CORE",
                        "Описание проекта": "Очень длинное описание " * 120,
                        "Зачем и для чего делаем": "Ещё текст " * 80,
                    },
                ],
                totalShown=2,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    core_slides = [
        prs.slides[index]
        for index in range(FIXED_SLIDE_COUNT, len(prs.slides))
        if any(
            shape.has_text_frame and shape.text_frame.text == "Продуктовый офис: CORE"
            for shape in prs.slides[index].shapes
        )
    ]
    assert len(core_slides) >= 2


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_fills_notes_for_legacy_sheet(mock_load) -> None:
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=["Дата запуска", "Проект", "Описание проекта", "Зачем и для чего делаем"],
                rows=[
                    {
                        "Дата запуска": "01.06",
                        "Проект": "CORE",
                        "Описание проекта": "Полный статус",
                        "Зачем и для чего делаем": "Полный зачем",
                    },
                ],
                totalShown=1,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    notes = content_slide.notes_slide.notes_text_frame.text
    assert "CORE" in notes
    assert "Полный статус" in notes
    assert "Полный зачем" in notes


@patch("app.product_status_presentation.load_b2b_product_status")
def test_generate_presentation_embeds_grid_table_with_borders(mock_load) -> None:
    import base64
    import json

    payload = {
        "text": "Договор с РТК",
        "table": {
            "rows": 2,
            "cols": 2,
            "cells": [["A1", "A2"], ["B1", "B2"]],
        },
    }
    token = (
        "<<tablejson:"
        + base64.b64encode(json.dumps(payload, ensure_ascii=False).encode("utf-8")).decode("ascii")
        + ">>"
    )
    mock_load.return_value = ProductStatusB2BOut(
        title="Статус продукта B2B",
        sheets=[
            ProductStatusSheetOut(
                gid="0",
                name="Продуктовый офис: CORE",
                columns=[
                    "Идет в презентацию",
                    "Дата запуска",
                    "Проект",
                    "Описание проекта и статус",
                    "Зачем и для чего делаем",
                ],
                rows=[
                    {
                        "Идет в презентацию": "Да",
                        "Дата запуска": "09.06",
                        "Проект": "РТК",
                        "Описание проекта и статус": token,
                        "Зачем и для чего делаем": "Контракт",
                    }
                ],
                totalShown=1,
            )
        ],
    )

    content, _ = generate_b2b_product_status_presentation()
    prs = Presentation(io.BytesIO(content))
    content_slide = prs.slides[FIXED_SLIDE_COUNT]
    tables = [shape.table for shape in content_slide.shapes if shape.has_table]
    assert len(tables) >= 2
    main_table = max(tables, key=lambda item: len(item.rows) * len(item.columns))
    nested = min(tables, key=lambda item: len(item.rows) * len(item.columns))
    assert "Договор с РТК" in main_table.rows[1].cells[2].text
    assert nested.rows[0].cells[0].text.strip() == "A1"
    assert nested.rows[1].cells[1].text.strip() == "B2"
    xml = content_slide.part.blob.decode()
    assert xml.count("<a:tbl>") >= 2
    assert 'srgbClr val="7F7F7F"' in xml
